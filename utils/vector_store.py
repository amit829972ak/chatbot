import os
import json
import re
from typing import Dict, List, Any, Optional
from datetime import datetime
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

# Document processing constants
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

try:
    # Try to use advanced embeddings with sentence transformers
    from sentence_transformers import SentenceTransformer
    import numpy as np
    from sklearn.metrics.pairwise import cosine_similarity
    
    class AdvancedVectorStore:
        """
        Advanced vector store implementation using sentence transformers for semantic search.
        Falls back to keyword search if embeddings fail.
        """
        def __init__(self, directory: str = "advanced_vector_store"):
            self.directory = directory
            self.index_file = os.path.join(directory, "index.json")
            self.embeddings_file = os.path.join(directory, "embeddings.npy")
            self.versions_file = os.path.join(directory, "versions.json")
            
            # Create directory if it doesn't exist
            os.makedirs(directory, exist_ok=True)
            
            # Initialize sentence transformer model
            try:
                self.model = SentenceTransformer('all-MiniLM-L6-v2')
                self.use_embeddings = True
            except Exception as e:
                print(f"Warning: Could not load sentence transformer model: {e}")
                self.use_embeddings = False
            
            # Load existing data
            self.content_index = self._load_index()
            self.embeddings_matrix = self._load_embeddings()
            self.versions = self._load_versions()
            
        def _load_index(self) -> Dict[str, Any]:
            """Load the content index from disk."""
            if os.path.exists(self.index_file):
                try:
                    with open(self.index_file, 'r', encoding='utf-8') as f:
                        return json.load(f)
                except Exception as e:
                    print(f"Error loading index: {e}")
            return {"documents": {}, "chunks": []}
        
        def _save_index(self) -> None:
            """Save the content index to disk."""
            try:
                with open(self.index_file, 'w', encoding='utf-8') as f:
                    json.dump(self.content_index, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"Error saving index: {e}")
                
        def _load_embeddings(self) -> Optional[np.ndarray]:
            """Load embeddings matrix from disk."""
            if os.path.exists(self.embeddings_file) and self.use_embeddings:
                try:
                    return np.load(self.embeddings_file)
                except Exception as e:
                    print(f"Error loading embeddings: {e}")
            return None
            
        def _save_embeddings(self) -> None:
            """Save embeddings matrix to disk."""
            if self.embeddings_matrix is not None and self.use_embeddings:
                try:
                    np.save(self.embeddings_file, self.embeddings_matrix)
                except Exception as e:
                    print(f"Error saving embeddings: {e}")
                    
        def _load_versions(self) -> Dict[str, Any]:
            """Load document versions from disk."""
            if os.path.exists(self.versions_file):
                try:
                    with open(self.versions_file, 'r', encoding='utf-8') as f:
                        return json.load(f)
                except Exception as e:
                    print(f"Error loading versions: {e}")
            return {}
            
        def _save_versions(self) -> None:
            """Save document versions to disk."""
            try:
                with open(self.versions_file, 'w', encoding='utf-8') as f:
                    json.dump(self.versions, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"Error saving versions: {e}")
        
        def add_document(self, document_path: str, content: str) -> None:
            """Add document content to the index with versioning."""
            # Split content into chunks
            chunks = self._split_text(content)
            
            # Check if document already exists for versioning
            doc_name = os.path.basename(document_path)
            is_update = doc_name in self.versions
            
            if is_update:
                # Remove old chunks for this document
                self.remove_document(document_path)
            
            # Add new chunks
            start_idx = len(self.content_index["chunks"])
            
            for i, chunk in enumerate(chunks):
                chunk_data = {
                    "content": chunk,
                    "document": document_path,
                    "chunk_id": start_idx + i,
                    "timestamp": datetime.now().isoformat()
                }
                self.content_index["chunks"].append(chunk_data)
            
            # Update document info
            self.content_index["documents"][document_path] = {
                "chunk_count": len(chunks),
                "total_length": len(content),
                "timestamp": datetime.now().isoformat(),
                "start_chunk_idx": start_idx,
                "end_chunk_idx": start_idx + len(chunks) - 1
            }
            
            # Update version tracking
            if doc_name not in self.versions:
                self.versions[doc_name] = {
                    "current_version": 1,
                    "versions": []
                }
            else:
                self.versions[doc_name]["current_version"] += 1
            
            # Add version entry
            version_entry = {
                "version": self.versions[doc_name]["current_version"],
                "timestamp": datetime.now().isoformat(),
                "size": len(content),
                "chunks": len(chunks)
            }
            self.versions[doc_name]["versions"].append(version_entry)
            
            # Generate embeddings for new chunks if possible
            if self.use_embeddings and chunks:
                try:
                    new_embeddings = self.model.encode(chunks)
                    
                    if self.embeddings_matrix is None:
                        self.embeddings_matrix = new_embeddings
                    else:
                        self.embeddings_matrix = np.vstack([self.embeddings_matrix, new_embeddings])
                    
                    self._save_embeddings()
                except Exception as e:
                    print(f"Error generating embeddings: {e}")
                    self.use_embeddings = False
            
            # Save updates
            self._save_index()
            self._save_versions()
        
        def _split_text(self, text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
            """Split text into overlapping chunks."""
            if len(text) <= chunk_size:
                return [text]
            
            chunks = []
            start = 0
            
            while start < len(text):
                end = start + chunk_size
                
                # Try to break at sentence boundary
                if end < len(text):
                    # Look for sentence ending
                    sentence_end = text.rfind('.', start, end)
                    if sentence_end > start + chunk_size // 2:
                        end = sentence_end + 1
                    else:
                        # Look for word boundary
                        word_end = text.rfind(' ', start, end)
                        if word_end > start:
                            end = word_end
                
                chunk = text[start:end].strip()
                if chunk:
                    chunks.append(chunk)
                
                start = max(start + chunk_size - overlap, end)
                
                if start >= len(text):
                    break
            
            return chunks
        
        def search(self, query: str, top_k: int = 3) -> List[str]:
            """
            Search for relevant chunks using embeddings or keyword search.
            Returns the top k most relevant chunks.
            """
            if not self.content_index["chunks"]:
                return []
            
            if self.use_embeddings and self.embeddings_matrix is not None:
                try:
                    return self._semantic_search(query, top_k)
                except Exception as e:
                    print(f"Semantic search failed, falling back to keyword search: {e}")
            
            return self._keyword_search(query, top_k)
        
        def _semantic_search(self, query: str, top_k: int) -> List[str]:
            """Perform semantic search using embeddings."""
            # Generate query embedding
            query_embedding = self.model.encode([query])
            
            # Calculate similarities
            similarities = cosine_similarity(query_embedding, self.embeddings_matrix)[0]
            
            # Get top k indices
            top_indices = np.argsort(similarities)[::-1][:top_k]
            
            # Return corresponding chunks
            results = []
            for idx in top_indices:
                if idx < len(self.content_index["chunks"]):
                    chunk = self.content_index["chunks"][idx]
                    results.append(f"[{os.path.basename(chunk['document'])}] {chunk['content']}")
            
            return results
        
        def _keyword_search(self, query: str, top_k: int) -> List[str]:
            """
            Fallback keyword search functionality.
            Returns the top k chunks that contain the most terms from the query.
            """
            query_terms = set(re.findall(r'\b\w+\b', query.lower()))
            
            if not query_terms:
                return []
            
            scored_chunks = []
            
            for chunk in self.content_index["chunks"]:
                content_lower = chunk["content"].lower()
                content_terms = set(re.findall(r'\b\w+\b', content_lower))
                
                # Calculate match score
                matches = query_terms.intersection(content_terms)
                if matches:
                    score = len(matches) / len(query_terms)
                    scored_chunks.append((score, chunk))
            
            # Sort by score and return top k
            scored_chunks.sort(key=lambda x: x[0], reverse=True)
            
            results = []
            for score, chunk in scored_chunks[:top_k]:
                results.append(f"[{os.path.basename(chunk['document'])}] {chunk['content']}")
            
            return results
        
        def remove_document(self, document_path: str) -> None:
            """Remove a document and its chunks from the index."""
            if document_path not in self.content_index["documents"]:
                return
            
            doc_info = self.content_index["documents"][document_path]
            start_idx = doc_info["start_chunk_idx"]
            end_idx = doc_info["end_chunk_idx"]
            
            # Remove chunks
            chunks_to_remove = end_idx - start_idx + 1
            del self.content_index["chunks"][start_idx:end_idx + 1]
            
            # Remove embeddings if they exist
            if self.embeddings_matrix is not None and self.use_embeddings:
                try:
                    mask = np.ones(len(self.embeddings_matrix), dtype=bool)
                    mask[start_idx:end_idx + 1] = False
                    self.embeddings_matrix = self.embeddings_matrix[mask]
                    self._save_embeddings()
                except Exception as e:
                    print(f"Error removing embeddings: {e}")
            
            # Update chunk indices for remaining documents
            for doc_path, info in self.content_index["documents"].items():
                if info["start_chunk_idx"] > end_idx:
                    info["start_chunk_idx"] -= chunks_to_remove
                    info["end_chunk_idx"] -= chunks_to_remove
            
            # Remove document info
            del self.content_index["documents"][document_path]
            
            # Save updates
            self._save_index()

except ImportError:
    print("Advanced embeddings not available, using simple vector store")
    
    # Fallback to simple implementation if dependencies are missing
    class AdvancedVectorStore:
        """Fallback to simple implementation."""
        def __init__(self, directory: str = "simple_vector_store"):
            self.directory = directory
            self.index_file = os.path.join(directory, "index.json")
            self.versions_file = os.path.join(directory, "versions.json")
            
            os.makedirs(directory, exist_ok=True)
            
            self.content_index = self._load_index()
            self.versions = self._load_versions()
            
        def _load_index(self) -> Dict[str, Any]:
            if os.path.exists(self.index_file):
                try:
                    with open(self.index_file, 'r', encoding='utf-8') as f:
                        return json.load(f)
                except:
                    pass
            return {"documents": {}, "chunks": []}
        
        def _save_index(self) -> None:
            try:
                with open(self.index_file, 'w', encoding='utf-8') as f:
                    json.dump(self.content_index, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"Error saving index: {e}")
                
        def _load_versions(self) -> Dict[str, Any]:
            if os.path.exists(self.versions_file):
                try:
                    with open(self.versions_file, 'r', encoding='utf-8') as f:
                        return json.load(f)
                except:
                    pass
            return {}
            
        def _save_versions(self) -> None:
            try:
                with open(self.versions_file, 'w', encoding='utf-8') as f:
                    json.dump(self.versions, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"Error saving versions: {e}")
        
        def add_document(self, document_path: str, content: str) -> None:
            chunks = self._split_text(content)
            
            doc_name = os.path.basename(document_path)
            is_update = doc_name in self.versions
            
            if is_update:
                self.remove_document(document_path)
            
            start_idx = len(self.content_index["chunks"])
            
            for i, chunk in enumerate(chunks):
                chunk_data = {
                    "content": chunk,
                    "document": document_path,
                    "chunk_id": start_idx + i,
                    "timestamp": datetime.now().isoformat()
                }
                self.content_index["chunks"].append(chunk_data)
            
            self.content_index["documents"][document_path] = {
                "chunk_count": len(chunks),
                "total_length": len(content),
                "timestamp": datetime.now().isoformat(),
                "start_chunk_idx": start_idx,
                "end_chunk_idx": start_idx + len(chunks) - 1
            }
            
            if doc_name not in self.versions:
                self.versions[doc_name] = {"current_version": 1, "versions": []}
            else:
                self.versions[doc_name]["current_version"] += 1
            
            version_entry = {
                "version": self.versions[doc_name]["current_version"],
                "timestamp": datetime.now().isoformat(),
                "size": len(content),
                "chunks": len(chunks)
            }
            self.versions[doc_name]["versions"].append(version_entry)
            
            self._save_index()
            self._save_versions()
        
        def _split_text(self, text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
            if len(text) <= chunk_size:
                return [text]
            
            chunks = []
            start = 0
            
            while start < len(text):
                end = start + chunk_size
                
                if end < len(text):
                    sentence_end = text.rfind('.', start, end)
                    if sentence_end > start + chunk_size // 2:
                        end = sentence_end + 1
                    else:
                        word_end = text.rfind(' ', start, end)
                        if word_end > start:
                            end = word_end
                
                chunk = text[start:end].strip()
                if chunk:
                    chunks.append(chunk)
                
                start = max(start + chunk_size - overlap, end)
                
                if start >= len(text):
                    break
            
            return chunks
        
        def search(self, query: str, top_k: int = 3) -> List[str]:
            if not self.content_index["chunks"]:
                return []
            
            query_terms = set(re.findall(r'\b\w+\b', query.lower()))
            
            if not query_terms:
                return []
            
            scored_chunks = []
            
            for chunk in self.content_index["chunks"]:
                content_lower = chunk["content"].lower()
                content_terms = set(re.findall(r'\b\w+\b', content_lower))
                
                matches = query_terms.intersection(content_terms)
                if matches:
                    score = len(matches) / len(query_terms)
                    scored_chunks.append((score, chunk))
            
            scored_chunks.sort(key=lambda x: x[0], reverse=True)
            
            results = []
            for score, chunk in scored_chunks[:top_k]:
                results.append(f"[{os.path.basename(chunk['document'])}] {chunk['content']}")
            
            return results
        
        def remove_document(self, document_path: str) -> None:
            if document_path not in self.content_index["documents"]:
                return
            
            doc_info = self.content_index["documents"][document_path]
            start_idx = doc_info["start_chunk_idx"]
            end_idx = doc_info["end_chunk_idx"]
            
            chunks_to_remove = end_idx - start_idx + 1
            del self.content_index["chunks"][start_idx:end_idx + 1]
            
            for doc_path, info in self.content_index["documents"].items():
                if info["start_chunk_idx"] > end_idx:
                    info["start_chunk_idx"] -= chunks_to_remove
                    info["end_chunk_idx"] -= chunks_to_remove
            
            del self.content_index["documents"][document_path]
            self._save_index()


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from an Excel file."""
    try:
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"\n=== Sheet: {sheet_name} ===\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    text += " | ".join(row_text) + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        text = ""
        
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n=== Slide {i} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # Extract text from tables
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            text += " | ".join(row_text) + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text from TXT: {e}")
            return ""
    except Exception as e:
        print(f"Error extracting text from TXT: {e}")
        return ""

def extract_text_from_file(file_path: str) -> str:
    """Extract text from a file based on its extension."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return ""

def update_index_from_file(file_path: str) -> bool:
    """
    Process a file and update the vector store index.
    
    Args:
        file_path: Path to the file to process
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        # Extract text from the file
        text_content = extract_text_from_file(file_path)
        
        if not text_content:
            print(f"No text content extracted from {file_path}")
            return False
        
        # Create vector store instance
        vector_store = AdvancedVectorStore()
        
        # Add document to the vector store
        vector_store.add_document(file_path, text_content)
        
        return True
        
    except Exception as e:
        print(f"Error updating index from file: {e}")
        return False

def search_index(query: str, top_k: int = 3, specific_docs: Optional[List[str]] = None) -> List[str]:
    """
    Search the vector store index for relevant document chunks.
    
    Args:
        query: Query string to search for
        top_k: Number of results to return
        specific_docs: Optional list of specific document paths to search within
        
    Returns:
        List of relevant document chunks
    """
    try:
        # Create vector store instance
        vector_store = AdvancedVectorStore()
        
        # Search for relevant chunks
        results = vector_store.search(query, top_k)
        
        # Filter by specific documents if requested
        if specific_docs:
            filtered_results = []
            for result in results:
                for doc in specific_docs:
                    if f"[{doc}]" in result:
                        filtered_results.append(result)
                        break
            return filtered_results
        
        return results
        
    except Exception as e:
        print(f"Error searching index: {e}")
        return []

def remove_from_index(file_path: str) -> bool:
    """
    Remove a file from the vector store index.
    
    Args:
        file_path: Path to the file to remove
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        # Create vector store instance
        vector_store = AdvancedVectorStore()
        
        # Remove document from the vector store
        vector_store.remove_document(file_path)
        
        return True
        
    except Exception as e:
        print(f"Error removing file from index: {e}")
        return False
